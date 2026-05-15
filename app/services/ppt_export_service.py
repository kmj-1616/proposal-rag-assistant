"""PPT 내보내기 서비스: 섹션별 초안 데이터를 .pptx 파일로 변환한다.

슬라이드 구성:
- 표지 슬라이드: 프로젝트명
- 목차 슬라이드: 섹션 제목 목록
- 섹션 슬라이드: 섹션마다 제목 + 내용 (긴 내용은 자동 분할)
"""
from __future__ import annotations

import re
from io import BytesIO
from typing import Any

_MAX_LINES_PER_SLIDE = 12  # 슬라이드 한 장당 최대 표시 줄 수


def _strip_markup(text: str) -> str:
    """마크다운 인라인 기호와 헤딩 기호를 제거하고 순수 텍스트를 반환한다."""
    text = re.sub(r"^\s*#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


def _parse_content_lines(content: str) -> list[str]:
    """content 문자열을 슬라이드에 표시할 줄 목록으로 변환한다.

    처리 규칙:
    - ```...``` 코드 블록: 스킵
    - # 헤딩: 「▶ 제목」 형태로 변환
    - - / * 글머리: 「• 내용」으로 변환
    - 1. / 1) 번호 목록: 「1. 내용」으로 유지
    - 빈 줄: 스킵
    - 일반 텍스트: 그대로 포함
    """
    lines: list[str] = []
    in_code_block = False

    for raw in content.splitlines():
        if raw.strip().startswith("```"):
            in_code_block = not in_code_block
            continue
        if in_code_block:
            continue

        stripped = raw.rstrip()
        if not stripped:
            continue

        # 헤딩
        m = re.match(r"^\s*#{1,6}\s+(.*)", stripped)
        if m:
            lines.append(f"▶ {_strip_markup(m.group(1).strip())}")
            continue

        # 글머리 기호
        m = re.match(r"^\s*[-*]\s+(.*)", stripped)
        if m:
            lines.append(f"• {_strip_markup(m.group(1).strip())}")
            continue

        # 번호 목록
        m = re.match(r"^\s*(\d+[.)]\s+.*)", stripped)
        if m:
            lines.append(_strip_markup(m.group(1).strip()))
            continue

        lines.append(_strip_markup(stripped))

    return lines


def _chunk_lines(lines: list[str], max_lines: int) -> list[list[str]]:
    """줄 목록을 max_lines 단위로 분할한다."""
    if not lines:
        return [[]]
    return [lines[i : i + max_lines] for i in range(0, len(lines), max_lines)]


def _set_title_slide(slide: Any, title_text: str, subtitle_text: str = "") -> None:
    """표지 슬라이드의 제목과 부제목을 설정한다."""
    tf = slide.shapes.title
    if tf:
        tf.text = title_text

    placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
    if placeholders and subtitle_text:
        placeholders[0].text = subtitle_text


def _add_content_slide(
    prs: Any,
    layout: Any,
    title_text: str,
    lines: list[str],
    page_suffix: str = "",
) -> None:
    """제목 + 내용 슬라이드를 프레젠테이션에 추가한다."""
    from pptx.util import Pt

    slide = prs.slides.add_slide(layout)

    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title_text + page_suffix

    # 본문 placeholder (idx=1)
    body_ph = next(
        (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
        None,
    )
    if body_ph is None:
        return

    tf = body_ph.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        if para.runs:
            para.runs[0].font.size = Pt(16)


def build_pptx(project_name: str | None, sections: list[dict[str, Any]]) -> bytes:
    """섹션 목록을 받아 PowerPoint 파일 바이트를 반환한다.

    Args:
        project_name: 표지 제목으로 사용할 프로젝트명. None이면 기본 제목 사용.
        sections: [{"title": str, "content": str, "priority": int, ...}] 형태의 목록.

    Returns:
        .pptx 파일 바이트.
    """
    from pptx import Presentation

    prs = Presentation()

    title_layout = prs.slide_layouts[0]   # 표지 레이아웃
    content_layout = prs.slide_layouts[1]  # 제목+내용 레이아웃

    display_title = project_name or "제안서 초안"

    # 표지 슬라이드
    cover = prs.slides.add_slide(title_layout)
    _set_title_slide(cover, display_title, "제안서 자동 생성 시스템")

    # 목차 슬라이드
    toc_slide = prs.slides.add_slide(content_layout)
    toc_title = toc_slide.shapes.title
    if toc_title:
        toc_title.text = "목차"
    toc_body = next(
        (ph for ph in toc_slide.placeholders if ph.placeholder_format.idx == 1),
        None,
    )
    if toc_body:
        tf = toc_body.text_frame
        tf.clear()
        for i, section in enumerate(sections, start=1):
            para = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            para.text = f"{i}. {section.get('title', '')}"

    # 섹션 슬라이드
    for section in sections:
        sec_title = section.get("title", "")
        content = section.get("content", "")

        lines = _parse_content_lines(content)
        chunks = _chunk_lines(lines, _MAX_LINES_PER_SLIDE)

        for page_idx, chunk in enumerate(chunks):
            suffix = f" ({page_idx + 1}/{len(chunks)})" if len(chunks) > 1 else ""
            _add_content_slide(prs, content_layout, sec_title, chunk, suffix)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
