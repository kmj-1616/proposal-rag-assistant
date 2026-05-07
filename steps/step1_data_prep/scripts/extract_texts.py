"""Extract plain text from local proposal files (.pptx/.pdf)."""
from __future__ import annotations

import json
import re
import sys
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path

from path_config import ensure_dir, step1_dir


@dataclass
class FileReport:
    source: str
    ok: bool
    chars: int
    error: str | None
    output: str | None


def _safe_stem(name: str) -> str:
    stem = Path(name).stem
    stem = re.sub(r'[<>:"/\\|?*]', "_", stem)
    return stem or "unnamed"


def _pptx_text(path: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    blocks: list[str] = []

    def walk_shapes(shapes, depth: int = 0) -> None:
        if depth > 30:
            return
        for shape in shapes:
            st = getattr(shape, "shape_type", None)
            if st == MSO_SHAPE_TYPE.GROUP:
                walk_shapes(shape.shapes, depth + 1)
                continue
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    txt = (p.text or "").strip()
                    if txt:
                        blocks.append(txt)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = []
                    for cell in row.cells:
                        ct = (cell.text or "").strip().replace("\n", " ")
                        if ct:
                            cells.append(ct)
                    if cells:
                        blocks.append("\t".join(cells))

    for i, slide in enumerate(prs.slides, 1):
        blocks.append(f"\n--- Slide {i} ---\n")
        walk_shapes(slide.shapes)

    return "\n".join(blocks).strip() + "\n"


def _pdf_text_pymupdf(path: Path) -> str:
    import fitz

    doc = fitz.open(path)
    try:
        parts: list[str] = []
        for i in range(len(doc)):
            page = doc.load_page(i)
            t = page.get_text("text") or ""
            parts.append(f"\n--- PDF Page {i + 1} ---\n")
            parts.append(t)
        return "\n".join(parts).strip() + "\n"
    finally:
        doc.close()


def _pdf_text_pdfplumber(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts.append(f"\n--- PDF Page {i} ---\n")
            parts.append(page.extract_text() or "")
    return "\n".join(parts).strip() + "\n"


def _pdf_text(path: Path) -> str:
    try:
        text = _pdf_text_pymupdf(path)
        if len(re.sub(r"\s+", "", text)) < 20:
            return _pdf_text_pdfplumber(path)
        return text
    except Exception:
        return _pdf_text_pdfplumber(path)


def main() -> int:
    src_dir = step1_dir("raw_proposals")
    out_dir = ensure_dir(step1_dir("extracted_texts"))

    if not src_dir.is_dir():
        print(f"Missing raw input directory: {src_dir}", file=sys.stderr)
        return 1

    files = sorted(
        p for p in src_dir.iterdir() if p.is_file() and p.suffix.lower() in (".pptx", ".pdf")
    )
    if not files:
        print(f"No supported input files found: {src_dir}", file=sys.stderr)
        return 1

    reports: list[FileReport] = []
    for path in files:
        out_path = out_dir / f"{_safe_stem(path.name)}.txt"
        error: str | None = None
        text = ""
        try:
            if path.suffix.lower() == ".pptx":
                text = _pptx_text(path)
            else:
                text = _pdf_text(path)
            out_path.write_text(text, encoding="utf-8", newline="\n")
        except Exception as exc:  # noqa: BLE001
            error = f"{type(exc).__name__}: {exc}"
            out_path.write_text(f"[EXTRACTION_FAILED] {error}\n", encoding="utf-8", newline="\n")

        reports.append(
            FileReport(
                source=path.name,
                ok=error is None,
                chars=len(text) if error is None else 0,
                error=error,
                output=out_path.name,
            )
        )

    summary = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_dir": str(src_dir),
        "output_dir": str(out_dir),
        "files": [asdict(r) for r in reports],
    }
    (out_dir / "_extraction_summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )

    ok_count = sum(1 for r in reports if r.ok)
    print(f"Extraction done: {ok_count}/{len(reports)} succeeded")
    return 0 if ok_count == len(reports) else 2


if __name__ == "__main__":
    raise SystemExit(main())
